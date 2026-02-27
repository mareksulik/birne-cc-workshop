#!/usr/bin/env python3
"""
Build a PPTX from screenshot PNGs + speaker notes extracted from index.html.

Usage: python3 build-pptx.py

Expects:
  - exports/slides/slide-001.png, slide-002.png, ...
  - index.html in the same directory

Produces:
  - workshop-prezentacia.pptx
"""

from __future__ import annotations

import glob
import os
import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches

SCRIPT_DIR = Path(__file__).parent
SLIDES_DIR = SCRIPT_DIR / "exports" / "slides"
INDEX_HTML = SCRIPT_DIR / "index.html"
OUTPUT_FILE = SCRIPT_DIR / "workshop-prezentacia.pptx"

# 16:9 slide size matching 1280×720
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)


def extract_notes(html_path: Path) -> list[str]:
    """Extract speaker notes from each top-level <section> in Reveal.js HTML."""
    with open(html_path, encoding="utf-8") as f:
        soup = BeautifulSoup(f.read(), "html.parser")

    slides_container = soup.select_one(".slides")
    if not slides_container:
        print("ERROR: Could not find .slides container in HTML", file=sys.stderr)
        sys.exit(1)

    notes = []
    # Top-level sections only (direct children of .slides)
    for section in slides_container.find_all("section", recursive=False):
        aside = section.find("aside", class_="notes")
        if aside:
            # Clean up whitespace but preserve line breaks
            text = aside.get_text(separator="\n").strip()
            notes.append(text)
        else:
            notes.append("")

    return notes


def get_slide_images(slides_dir: Path) -> list[Path]:
    """Return sorted list of slide PNG paths."""
    files = sorted(slides_dir.glob("slide-*.png"))
    return files


def build_pptx(images: list[Path], notes: list[str], output: Path):
    """Assemble PPTX with screenshot slides and speaker notes."""
    prs = Presentation()

    # Set slide dimensions to 16:9
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank

    for idx, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)

        # Add image as full-bleed background
        slide.shapes.add_picture(
            str(img_path),
            left=0,
            top=0,
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )

        # Add speaker notes if available
        if idx < len(notes) and notes[idx]:
            slide.notes_slide.notes_text_frame.text = notes[idx]

    prs.save(str(output))
    print(f"Saved {output} with {len(images)} slides")


def main():
    # Extract notes
    if not INDEX_HTML.exists():
        print(f"ERROR: {INDEX_HTML} not found", file=sys.stderr)
        sys.exit(1)

    notes = extract_notes(INDEX_HTML)
    print(f"Extracted notes from {len(notes)} sections")

    # Get slide images
    images = get_slide_images(SLIDES_DIR)
    if not images:
        print(f"ERROR: No slide images found in {SLIDES_DIR}", file=sys.stderr)
        print("Run 'node export-slides.js' first to generate screenshots.", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(images)} slide images")

    if len(images) != len(notes):
        print(
            f"WARNING: {len(images)} images but {len(notes)} sections in HTML. "
            f"Notes may not align perfectly.",
            file=sys.stderr,
        )

    # Build PPTX
    build_pptx(images, notes, OUTPUT_FILE)


if __name__ == "__main__":
    main()
